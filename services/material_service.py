from typing import Optional, Tuple
from datetime import datetime
import io
import os
import uuid

from models.schemas import Material, MaterialType
from storage.db import material_db   # ⭐ USE SQLITE NOW


# -------------------------------
# DETECT FILE TYPE
# -------------------------------
def detect_material_type(filename: str, content_type: Optional[str] = None) -> MaterialType:
    name = filename.lower()
    if name.endswith('.pdf'):
        return MaterialType.PDF
    if name.endswith('.docx'):
        return MaterialType.DOCX
    if name.endswith('.pptx'):
        return MaterialType.PPTX
    if name.endswith('.txt'):
        return MaterialType.TXT
    if name.endswith('.md'):
        return MaterialType.MARKDOWN
    return MaterialType.UNKNOWN


# -------------------------------
# PDF / DOCX / PPTX EXTRACTORS
# -------------------------------
def extract_text_from_pdf(file_bytes: bytes) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        pages = [p.extract_text() or "" for p in reader.pages]
        return "\n\n".join(pages)
    except Exception as e:
        return f"[Error extracting PDF: {str(e)}]"


def extract_text_from_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(parts)
    except Exception as e:
        return f"[Error extracting DOCX: {str(e)}]"


def extract_text_from_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            items = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    items.append(shape.text)
            slides.append("\n".join(items))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[Error extracting PPTX: {str(e)}]"


# -------------------------------
# EXTRACT TEXT BASED ON TYPE
# -------------------------------
def extract_text(file_bytes: bytes, material_type: MaterialType) -> str:
    if material_type == MaterialType.PDF:
        return extract_text_from_pdf(file_bytes)
    if material_type == MaterialType.DOCX:
        return extract_text_from_docx(file_bytes)
    if material_type == MaterialType.PPTX:
        return extract_text_from_pptx(file_bytes)
    if material_type in (MaterialType.TXT, MaterialType.MARKDOWN):
        try:
            return file_bytes.decode("utf-8")
        except:
            return "[Error decoding text file]"
    return "[Unknown file format]"


# -----------------------------------------------------
# PROCESS AND STORE FILE → SAVE TO SQLITE
# -----------------------------------------------------
def process_and_store_file(
    filename: str,
    file_bytes: bytes,
    content_type: Optional[str] = None,
    custom_title: Optional[str] = None,
) -> Tuple[str, Material]:

    material_type = detect_material_type(filename)
    content = extract_text(file_bytes, material_type)

    material_id = str(uuid.uuid4())[:8]
    title = custom_title or os.path.splitext(filename)[0]

    material = Material(
        id=material_id,
        title=title,
        type=material_type,
        content=content,
        created_at=datetime.now().isoformat(),
        metadata_info=f"filename:{filename},size:{len(file_bytes)}"
    )

    material_db.save(material)   # ⭐ SAVE TO SQLITE

    return material_id, material


# -----------------------------------------------------
# PROCESS AND STORE RAW TEXT → SAVE TO SQLITE
# -----------------------------------------------------
def process_and_store_text(text: str, title: str = "Untitled Notes") -> Tuple[str, Material]:

    material_id = str(uuid.uuid4())[:8]

    material = Material(
        id=material_id,
        title=title,
        type=MaterialType.TXT,
        content=text,
        created_at=datetime.now().isoformat(),
        metadata_info="source:direct_text"
    )

    material_db.save(material)   # ⭐ SAVE TO SQLITE

    return material_id, material


# -----------------------------------------------------
# READ CONTENT FROM SQLITE
# -----------------------------------------------------
def get_material_content(material_id: str) -> Optional[str]:
    row = material_db.get(material_id)
    return row.content if row else None


def get_all_materials_content(material_ids: Optional[list] = None) -> str:
    if material_ids:
        texts = []
        for mid in material_ids:
            row = material_db.get(mid)
            if row:
                texts.append(row.content)
        return "\n\n".join(texts)

    # all materials
    rows = material_db.list()
    return "\n\n".join([r.content for r in rows]) if rows else "[No materials uploaded]"


